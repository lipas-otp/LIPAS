"""Bounded, dependency-optional document extraction and conversion helpers.

The Workbench owns path policy, staging, approvals, and evidence.  This module
only handles bytes that the Workbench has already authorized.  Optional parser
packages are imported lazily so the LIPAS core remains dependency-free.
"""
from __future__ import annotations

import csv
import html
import io
import json
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Mapping

__all__ = [
    "DocumentToolError",
    "MissingDocumentDependency",
    "UnsupportedDocumentFormat",
    "ConvertedDocument",
    "convert_document",
    "read_pdf_text",
]


class DocumentToolError(ValueError):
    """A bounded document operation could not be completed."""


class MissingDocumentDependency(DocumentToolError):
    """An optional parser required by the selected document format is absent."""


class UnsupportedDocumentFormat(DocumentToolError):
    """The requested source or target format is not supported."""


@dataclass(frozen=True, slots=True)
class ConvertedDocument:
    """Deterministic conversion output and non-sensitive extraction metadata."""

    content: bytes
    source_format: str
    target_format: str
    metadata: Mapping[str, Any]


_FORMAT_ALIASES = {
    "text": "txt",
    "txt": "txt",
    "markdown": "md",
    "md": "md",
    "htm": "html",
    "html": "html",
    "json": "json",
    "jsonl": "json",
    "csv": "csv",
    "pdf": "pdf",
    "docx": "docx",
    "xlsx": "xlsx",
    "pptx": "pptx",
}
_TEXT_FORMATS = frozenset({"txt", "md", "html"})
_TABULAR_FORMATS = frozenset({"csv", "xlsx"})
_MAX_PDF_PAGES = 200
_MAX_CELLS = 100_000


def _normalise_format(value: str, *, field: str) -> str:
    if not isinstance(value, str) or not value.strip():
        raise UnsupportedDocumentFormat(f"{field} format must be a non-empty string")
    selected = value.strip().lower().lstrip(".")
    try:
        return _FORMAT_ALIASES[selected]
    except KeyError as exc:
        available = ", ".join(sorted(set(_FORMAT_ALIASES.values())))
        raise UnsupportedDocumentFormat(
            f"unsupported {field} format {value!r}; supported: {available}",
        ) from exc


def _read_bytes(path: Path, max_bytes: int) -> bytes:
    if isinstance(max_bytes, bool) or not isinstance(max_bytes, int) or max_bytes < 1:
        raise DocumentToolError("max_bytes must be a positive integer")
    try:
        size = path.stat().st_size
    except OSError as exc:
        raise DocumentToolError(f"cannot inspect source file: {exc}") from exc
    if size > max_bytes:
        raise DocumentToolError(
            f"source file exceeds the {max_bytes} byte document limit",
        )
    try:
        return path.read_bytes()
    except OSError as exc:
        raise DocumentToolError(f"cannot read source file: {exc}") from exc


def _decode_text(data: bytes) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise DocumentToolError("text document is not valid UTF-8") from exc


def _html_to_text(value: str) -> str:
    # A deliberately conservative conversion: remove tags while preserving
    # visible text and line boundaries.  It does not execute scripts or load
    # external resources.
    value = re.sub(r"(?is)<(script|style)\b[^>]*>.*?</\1>", "", value)
    value = re.sub(r"(?i)<\s*br\s*/?\s*>", "\n", value)
    value = re.sub(r"(?i)</\s*(p|div|h[1-6]|li|tr)\s*>", "\n", value)
    value = re.sub(r"(?s)<[^>]+>", "", value)
    value = html.unescape(value)
    lines = [line.rstrip() for line in value.splitlines()]
    return "\n".join(lines).strip()


def _markdown_to_html(value: str) -> str:
    blocks: list[str] = []
    paragraph: list[str] = []
    in_code = False
    code_lines: list[str] = []

    def flush_paragraph() -> None:
        if paragraph:
            blocks.append(f"<p>{html.escape(' '.join(paragraph))}</p>")
            paragraph.clear()

    for raw in value.splitlines():
        line = raw.rstrip()
        if line.startswith("```"):
            if in_code:
                blocks.append(f"<pre><code>{html.escape(chr(10).join(code_lines))}</code></pre>")
                code_lines.clear()
            else:
                flush_paragraph()
            in_code = not in_code
            continue
        if in_code:
            code_lines.append(line)
            continue
        if not line.strip():
            flush_paragraph()
            continue
        heading = re.match(r"^(#{1,6})\s+(.*)$", line)
        if heading:
            flush_paragraph()
            level = len(heading.group(1))
            blocks.append(f"<h{level}>{html.escape(heading.group(2).strip())}</h{level}>")
            continue
        if line.lstrip().startswith(("- ", "* ")):
            flush_paragraph()
            item = line.lstrip()[2:].strip()
            if blocks and blocks[-1].startswith("<ul>"):
                blocks[-1] = blocks[-1][:-5] + f"<li>{html.escape(item)}</li></ul>"
            else:
                blocks.append(f"<ul><li>{html.escape(item)}</li></ul>")
            continue
        paragraph.append(line.strip())
    if in_code:
        blocks.append(f"<pre><code>{html.escape(chr(10).join(code_lines))}</code></pre>")
    flush_paragraph()
    return "<!doctype html>\n<html><body>\n" + "\n".join(blocks) + "\n</body></html>\n"


def _as_text(data: bytes, source_format: str) -> str:
    if source_format in _TEXT_FORMATS:
        decoded = _decode_text(data)
        return _html_to_text(decoded) if source_format == "html" else decoded
    if source_format == "pdf":
        return read_pdf_text_from_bytes(data)[0]
    if source_format == "docx":
        try:
            from docx import Document  # type: ignore[import-not-found]
        except ImportError as exc:
            raise MissingDocumentDependency(
                "DOCX support requires `pip install 'lipas[documents]'`",
            ) from exc
        try:
            document = Document(io.BytesIO(data))
            paragraphs = [paragraph.text for paragraph in document.paragraphs]
            for table in document.tables:
                paragraphs.extend(
                    "\t".join(cell.text for cell in row.cells)
                    for row in table.rows
                )
            return "\n".join(paragraphs).strip()
        except Exception as exc:
            raise DocumentToolError(f"could not parse DOCX: {exc}") from exc
    if source_format == "xlsx":
        return _xlsx_to_tsv(data)
    if source_format == "pptx":
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
        except ImportError as exc:
            raise MissingDocumentDependency(
                "PPTX support requires `pip install 'lipas[documents]'`",
            ) from exc
        try:
            presentation = Presentation(io.BytesIO(data))
            lines: list[str] = []
            for slide_number, slide in enumerate(presentation.slides, 1):
                lines.append(f"# Slide {slide_number}")
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        value = shape.text.strip()
                        if value:
                            lines.append(value)
                    if getattr(shape, "has_table", False):
                        for row in shape.table.rows:
                            lines.append("\t".join(cell.text for cell in row.cells))
            return "\n".join(lines).strip()
        except Exception as exc:
            raise DocumentToolError(f"could not parse PPTX: {exc}") from exc
    if source_format == "csv":
        decoded = _decode_text(data)
        rows = list(csv.reader(io.StringIO(decoded)))
        return "\n".join("\t".join(row) for row in rows)
    if source_format == "json":
        try:
            value = json.loads(_decode_text(data))
        except (json.JSONDecodeError, UnicodeDecodeError) as exc:
            raise DocumentToolError(f"could not parse JSON: {exc}") from exc
        return json.dumps(value, ensure_ascii=False, indent=2, sort_keys=True) + "\n"
    raise UnsupportedDocumentFormat(f"cannot extract text from {source_format!r}")


def _xlsx_to_tsv(data: bytes) -> str:
    try:
        import openpyxl  # type: ignore[import-not-found,import-untyped]
    except ImportError as exc:
        raise MissingDocumentDependency(
            "XLSX support requires `pip install 'lipas[documents]'`",
        ) from exc
    try:
        workbook = openpyxl.load_workbook(
            io.BytesIO(data), read_only=True, data_only=True,
        )
        try:
            rows: list[str] = []
            cells = 0
            for sheet in workbook.worksheets:
                rows.append(f"# Sheet: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    values = ["" if value is None else str(value) for value in row]
                    cells += len(values)
                    if cells > _MAX_CELLS:
                        raise DocumentToolError("XLSX exceeds the cell limit")
                    rows.append("\t".join(values))
            return "\n".join(rows).strip()
        finally:
            close = getattr(workbook, "close", None)
            if callable(close):
                close()
    except DocumentToolError:
        raise
    except Exception as exc:
        raise DocumentToolError(f"could not parse XLSX: {exc}") from exc


def read_pdf_text(
    path: Path,
    *,
    max_bytes: int = 20 * 1024 * 1024,
    max_pages: int = _MAX_PDF_PAGES,
    max_chars: int = 120_000,
) -> tuple[str, Mapping[str, Any]]:
    """Extract bounded text from a PDF without executing embedded content."""
    data = _read_bytes(path, max_bytes)
    text, metadata = read_pdf_text_from_bytes(
        data, max_pages=max_pages, max_chars=max_chars,
    )
    return text, metadata


def read_pdf_text_from_bytes(
    data: bytes,
    *,
    max_pages: int = _MAX_PDF_PAGES,
    max_chars: int = 120_000,
) -> tuple[str, Mapping[str, Any]]:
    if isinstance(max_pages, bool) or not isinstance(max_pages, int) or not 1 <= max_pages <= _MAX_PDF_PAGES:
        raise DocumentToolError(f"max_pages must be between 1 and {_MAX_PDF_PAGES}")
    if isinstance(max_chars, bool) or not isinstance(max_chars, int) or not 1 <= max_chars <= 1_000_000:
        raise DocumentToolError("max_chars must be between 1 and 1000000")
    try:
        from pypdf import PdfReader  # type: ignore[import-not-found]
    except ImportError as exc:
        raise MissingDocumentDependency(
            "PDF support requires `pip install 'lipas[documents]'`",
        ) from exc
    try:
        reader = PdfReader(io.BytesIO(data), strict=False)
        if reader.is_encrypted:
            raise DocumentToolError("encrypted PDFs require explicit user-side decryption")
        page_count = len(reader.pages)
        chunks: list[str] = []
        pages_read = 0
        hit_character_limit = False
        for index, page in enumerate(reader.pages):
            if index >= max_pages:
                break
            chunks.append(page.extract_text() or "")
            pages_read += 1
            if sum(len(chunk) for chunk in chunks) >= max_chars:
                hit_character_limit = True
                break
        text = "\n\n".join(chunks).strip()[:max_chars]
        return text, {
            "pages": page_count,
            "pages_read": pages_read,
            "truncated": hit_character_limit or pages_read < page_count,
            "needs_ocr": bool(page_count and not text),
        }
    except DocumentToolError:
        raise
    except Exception as exc:
        raise DocumentToolError(f"could not parse PDF: {exc}") from exc


def _json_document(text: str, source_format: str) -> bytes:
    return (
        json.dumps(
            {"source_format": source_format, "text": text},
            ensure_ascii=False,
            indent=2,
            sort_keys=True,
        )
        + "\n"
    ).encode("utf-8")


def convert_document(
    path: Path,
    *,
    target_format: str,
    max_bytes: int = 20 * 1024 * 1024,
    max_chars: int = 120_000,
) -> ConvertedDocument:
    """Convert one authorized document to a bounded, reviewable format."""
    source_format = _normalise_format(path.suffix, field="source")
    target = _normalise_format(target_format, field="target")
    if target == source_format:
        raise UnsupportedDocumentFormat("source and target formats are identical")
    data = _read_bytes(path, max_bytes)

    if source_format in _TABULAR_FORMATS and target == "csv":
        if source_format == "csv":
            raise UnsupportedDocumentFormat("source and target formats are identical")
        text = _xlsx_to_tsv(data)
        rows = [line.split("\t") for line in text.splitlines() if not line.startswith("# Sheet:")]
        csv_output = io.StringIO(newline="")
        csv.writer(csv_output, lineterminator="\n").writerows(rows)
        return ConvertedDocument(
            csv_output.getvalue().encode("utf-8"), source_format, target,
            {"rows": len(rows)},
        )

    text = _as_text(data, source_format)
    if len(text) > max_chars:
        text = text[:max_chars]
        truncated = True
    else:
        truncated = False
    output: bytes
    if target in {"txt", "md"}:
        output = text.encode("utf-8")
        if not output.endswith(b"\n"):
            output += b"\n"
    elif target == "html":
        output = (
            _markdown_to_html(text) if source_format == "md" else
            "<!doctype html>\n<html><body><pre>" + html.escape(text) + "</pre></body></html>\n"
        ).encode("utf-8")
    elif target == "json":
        output = _json_document(text, source_format)
    elif target == "csv" and source_format == "xlsx":
        # Kept above for clarity; this branch is unreachable after the
        # tabular fast path but makes the supported pair explicit.
        raise UnsupportedDocumentFormat("XLSX CSV conversion failed")
    elif target == "docx":
        try:
            from docx import Document  # type: ignore[import-not-found]
        except ImportError as exc:
            raise MissingDocumentDependency(
                "DOCX output requires `pip install 'lipas[documents]'`",
            ) from exc
        try:
            document = Document()
            for line in text.splitlines() or [""]:
                document.add_paragraph(line)
            stream = io.BytesIO()
            document.save(stream)
            output = stream.getvalue()
        except Exception as exc:
            raise DocumentToolError(f"could not create DOCX: {exc}") from exc
    elif target == "xlsx":
        try:
            import openpyxl  # type: ignore[import-not-found,import-untyped]
        except ImportError as exc:
            raise MissingDocumentDependency(
                "XLSX output requires `pip install 'lipas[documents]'`",
            ) from exc
        try:
            workbook = openpyxl.Workbook()
            sheet = workbook.active
            assert sheet is not None
            rows = (
                list(csv.reader(io.StringIO(_decode_text(data))))
                if source_format == "csv"
                else [line.split("\t") for line in text.splitlines()]
            )
            for row_number, row in enumerate(rows, 1):
                for column_number, value in enumerate(row, 1):
                    sheet.cell(row=row_number, column=column_number, value=value)
            stream = io.BytesIO()
            workbook.save(stream)
            close = getattr(workbook, "close", None)
            if callable(close):
                close()
            output = stream.getvalue()
        except Exception as exc:
            raise DocumentToolError(f"could not create XLSX: {exc}") from exc
    else:
        raise UnsupportedDocumentFormat(
            f"conversion {source_format!r} -> {target!r} is not supported",
        )
    return ConvertedDocument(
        output,
        source_format,
        target,
        {"source_bytes": len(data), "output_bytes": len(output), "truncated": truncated},
    )
